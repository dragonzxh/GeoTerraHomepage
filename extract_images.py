# extract_images.py
"""
从PPT中提取所有图片
"""

from pptx import Presentation
import os

def extract_images_from_ppt(ppt_path, output_dir="images"):
    """从PPT中提取所有图片"""
    try:
        prs = Presentation(ppt_path)
        
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
            print(f"已创建目录: {output_dir}")
        
        image_count = 0
        
        print(f"正在从PPT提取图片: {ppt_path}")
        print(f"总共有 {len(prs.slides)} 张幻灯片\n")
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_image_count = 0
            for shape in slide.shapes:
                # 方法1: 直接图片对象
                if hasattr(shape, "image"):
                    try:
                        image = shape.image
                        # 获取图片扩展名
                        ext = image.ext
                        # 生成文件名（幻灯片编号_图片序号）
                        filename = f"slide_{slide_num:02d}_img_{slide_image_count:02d}{ext}"
                        filepath = os.path.join(output_dir, filename)
                        
                        # 保存图片
                        with open(filepath, 'wb') as f:
                            f.write(image.blob)
                        
                        print(f"幻灯片 {slide_num}: 已保存图片 -> {filename}")
                        image_count += 1
                        slide_image_count += 1
                    except Exception as e:
                        print(f"幻灯片 {slide_num}: 提取图片时出错 - {str(e)}")
                
                # 方法2: 检查是否是图片占位符
                elif hasattr(shape, "shape_type"):
                    # 尝试从图片占位符提取
                    try:
                        if hasattr(shape, "image") or shape.shape_type == 13:  # 13 = MSO_SHAPE_TYPE.PICTURE
                            if hasattr(shape, "image"):
                                image = shape.image
                                ext = image.ext if hasattr(image, 'ext') else '.png'
                                filename = f"slide_{slide_num:02d}_img_{slide_image_count:02d}{ext}"
                                filepath = os.path.join(output_dir, filename)
                                
                                with open(filepath, 'wb') as f:
                                    f.write(image.blob)
                                
                                print(f"幻灯片 {slide_num}: 已保存图片（占位符）-> {filename}")
                                image_count += 1
                                slide_image_count += 1
                    except:
                        pass
                
                # 方法3: 尝试从所有形状中提取图片数据
                try:
                    # 检查形状是否有图片相关的属性
                    if hasattr(shape, "_element"):
                        element = shape._element
                        # 查找图片关系
                        if hasattr(element, "nvPicPr"):
                            # 尝试提取图片
                            pass
                except:
                    pass
        
        print(f"\n{'='*60}")
        print(f"提取完成！总共提取了 {image_count} 张图片")
        print(f"图片保存在: {os.path.abspath(output_dir)}")
        print(f"{'='*60}")
        
        return image_count
        
    except Exception as e:
        print(f"错误: {str(e)}")
        print("\n请确保已安装 python-pptx:")
        print("pip install python-pptx")
        return 0


if __name__ == "__main__":
    ppt_file = "GEOTERRA SDN BHD Company Profile.pptx"
    
    if not os.path.exists(ppt_file):
        print(f"错误: 找不到文件 {ppt_file}")
        print("请确保PPT文件在当前目录下")
    else:
        extract_images_from_ppt(ppt_file)

