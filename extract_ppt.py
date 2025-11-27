# extract_ppt.py
"""
PPT内容提取脚本
用于从PowerPoint文件中提取所有文本、表格和图片信息
"""

from pptx import Presentation
import json
import os

def extract_ppt_content(ppt_path, output_format='txt'):
    """
    提取PPT文件的所有内容
    
    参数:
        ppt_path: PPT文件路径
        output_format: 输出格式 ('txt', 'json', 'markdown')
    """
    try:
        # 读取PPT文件
        prs = Presentation(ppt_path)
        
        # 存储所有内容
        slides_data = []
        
        print(f"正在读取PPT文件: {ppt_path}")
        print(f"总共有 {len(prs.slides)} 张幻灯片\n")
        
        # 遍历所有幻灯片
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_info = {
                'slide_number': slide_num,
                'title': '',
                'text_content': [],
                'tables': [],
                'images': []
            }
            
            print(f"处理幻灯片 {slide_num}...")
            
            # 遍历幻灯片中的所有形状
            for shape in slide.shapes:
                # 提取文本内容
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    
                    # 判断是否是标题（通常在第一个文本框）
                    if not slide_info['title'] and len(text) < 100:
                        slide_info['title'] = text
                    else:
                        slide_info['text_content'].append(text)
                    
                    print(f"  - 文本: {text[:50]}...")
                
                # 提取表格内容
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    slide_info['tables'].append(table_data)
                    print(f"  - 表格: {len(table_data)} 行 x {len(table_data[0]) if table_data else 0} 列")
                
                # 提取图片信息
                if hasattr(shape, "image"):
                    image_info = {
                        'filename': shape.image.filename if hasattr(shape.image, 'filename') else 'unknown',
                        'width': shape.width,
                        'height': shape.height
                    }
                    slide_info['images'].append(image_info)
                    print(f"  - 图片: {image_info['filename']}")
            
            slides_data.append(slide_info)
            print()
        
        # 根据输出格式保存
        base_name = os.path.splitext(ppt_path)[0]
        
        if output_format == 'json':
            output_path = f"{base_name}_extracted.json"
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(slides_data, f, ensure_ascii=False, indent=2)
            print(f"内容已保存到: {output_path}")
        
        elif output_format == 'markdown':
            output_path = f"{base_name}_extracted.md"
            with open(output_path, 'w', encoding='utf-8') as f:
                for slide in slides_data:
                    f.write(f"# 幻灯片 {slide['slide_number']}\n\n")
                    if slide['title']:
                        f.write(f"## {slide['title']}\n\n")
                    for text in slide['text_content']:
                        f.write(f"{text}\n\n")
                    for table in slide['tables']:
                        f.write("| " + " | ".join(table[0]) + " |\n")
                        f.write("| " + " | ".join(["---"] * len(table[0])) + " |\n")
                        for row in table[1:]:
                            f.write("| " + " | ".join(row) + " |\n")
                        f.write("\n")
                    f.write("---\n\n")
            print(f"内容已保存到: {output_path}")
        
        else:  # txt格式
            output_path = f"{base_name}_extracted.txt"
            with open(output_path, 'w', encoding='utf-8') as f:
                for slide in slides_data:
                    f.write(f"{'='*60}\n")
                    f.write(f"幻灯片 {slide['slide_number']}\n")
                    f.write(f"{'='*60}\n\n")
                    if slide['title']:
                        f.write(f"标题: {slide['title']}\n\n")
                    for text in slide['text_content']:
                        f.write(f"{text}\n\n")
                    for table in slide['tables']:
                        f.write("表格内容:\n")
                        for row in table:
                            f.write(" | ".join(row) + "\n")
                        f.write("\n")
                    if slide['images']:
                        f.write("图片信息:\n")
                        for img in slide['images']:
                            f.write(f"  - {img['filename']}\n")
                    f.write("\n")
            print(f"内容已保存到: {output_path}")
        
        return slides_data
        
    except Exception as e:
        print(f"错误: {str(e)}")
        print("\n请确保已安装 python-pptx:")
        print("pip install python-pptx")
        return None


def main():
    """主函数"""
    ppt_file = "GEOTERRA SDN BHD Company Profile.pptx"
    
    if not os.path.exists(ppt_file):
        print(f"错误: 找不到文件 {ppt_file}")
        print("请确保PPT文件在当前目录下")
        return
    
    print("="*60)
    print("PPT内容提取工具")
    print("="*60)
    print()
    
    # 提取内容（可以选择 'txt', 'json', 'markdown'）
    content = extract_ppt_content(ppt_file, output_format='txt')
    
    if content:
        print("\n提取完成！")
        print(f"共提取 {len(content)} 张幻灯片的内容")


if __name__ == "__main__":
    main()

